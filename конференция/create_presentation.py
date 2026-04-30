#!/usr/bin/env python3
"""
Generate conference presentation:
"Разработка программного комплекса оркестрации"

Theoretical / survey talk — orchestration as a discipline,
with the PAS helper as one practical example.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Theme colors (GUAP palette) ─────────────────────────────────────
DK1 = RGBColor(0x24, 0x28, 0x34)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT1 = RGBColor(0x00, 0x5A, 0xAA)   # blue
ACCENT2 = RGBColor(0xE7, 0x0F, 0x47)   # red
ACCENT3 = RGBColor(0x00, 0xBE, 0xF3)   # cyan
ACCENT4 = RGBColor(0x92, 0x69, 0xC9)   # purple
ACCENT5 = RGBColor(0xFF, 0x64, 0x18)   # orange
ACCENT6 = RGBColor(0x00, 0x9A, 0x49)   # green
DARK_BLUE = RGBColor(0x00, 0x2C, 0x5F)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
GRAY = RGBColor(0x6B, 0x72, 0x80)

FONT = 'Arial'

LEFT_MARGIN = Emu(550863)
TOP_CONTENT = Emu(1100000)

# ── Helpers ──────────────────────────────────────────────────────────

def clear_slide(slide):
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, size=Pt(18), bold=False, color=DK1, align=PP_ALIGN.LEFT,
             font_name=FONT, italic=False):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    run.font.italic = italic
    return p

def add_paragraph(tf, text, size=Pt(18), bold=False, color=DK1, align=PP_ALIGN.LEFT,
                  space_before=Pt(6), font_name=FONT, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    run.font.italic = italic
    return p

def add_bullet(tf, text, size=Pt(15), color=DK1, level=0, bold=False):
    p = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT
    pPr = p._p.get_or_add_pPr()
    char = '●' if level == 0 else '○'
    for old in pPr.findall(qn('a:buChar')) + pPr.findall(qn('a:buNone')):
        pPr.remove(old)
    pPr.append(pPr.makeelement(qn('a:buChar'), {'char': char}))
    return p

def add_title(slide, text):
    txBox = add_textbox(slide, LEFT_MARGIN, Emu(407976), Emu(9608330), Emu(400000))
    set_text(txBox.text_frame, text, size=Pt(28), color=ACCENT1)
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color,
                     text="", text_color=WHITE, text_size=Pt(14), bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(10)
        tf.margin_right = Pt(10)
        tf.margin_top = Pt(6)
        tf.margin_bottom = Pt(6)
        set_text(tf, text, size=text_size, bold=bold, color=text_color, align=PP_ALIGN.CENTER)
    return shape

def add_rect(slide, left, top, width, height, fill_color,
             text="", text_color=WHITE, text_size=Pt(14), bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        set_text(tf, text, size=text_size, bold=bold, color=text_color, align=PP_ALIGN.CENTER)
    return shape

def add_circle(slide, left, top, size, fill_color, text="", text_color=WHITE, text_size=Pt(20)):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if text:
        set_text(shape.text_frame, text, size=text_size, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    return shape

def add_arrow(slide, x1, y1, x2, y2, color=ACCENT1, width=Pt(2)):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = width
    ln = conn.line._ln
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return conn

def add_chevron(slide, left, top, width, height, fill_color, text="", text_color=WHITE, text_size=Pt(13)):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(14)
        tf.margin_right = Pt(6)
        set_text(tf, text, size=text_size, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    return shape

# ═════════════════════════════════════════════════════════════════════
def main():
    template_path = os.path.join(os.path.dirname(__file__), 'guap_16-9_2023.pptx')
    prs = Presentation(template_path)

    # Remove all existing slides
    while len(prs.slides) > 0:
        elem = prs.slides._sldIdLst[0]
        for attr_name in elem.attrib:
            if 'id' in attr_name.lower() and attr_name != 'id':
                prs.part.drop_rel(elem.attrib[attr_name])
                break
        prs.slides._sldIdLst.remove(elem)

    layout = prs.slide_layouts[0]

    # ─── SLIDE 1: Title ─────────────────────────────────────────────
    slide = prs.slides.add_slide(layout)
    clear_slide(slide)

    txBox = add_textbox(slide, Emu(550864), Emu(2000000), Emu(9500000), Emu(1800000))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_text(tf, "Разработка программного комплекса", size=Pt(38), color=ACCENT1)
    add_paragraph(tf, "оркестрации", size=Pt(38), color=ACCENT1, space_before=Pt(0))

    txBox2 = add_textbox(slide, Emu(550863), Emu(4200000), Emu(7000000), Emu(1000000))
    tf2 = txBox2.text_frame
    set_text(tf2, "Соколов Илья Дмитриевич", size=Pt(18), color=DK1)
    add_paragraph(tf2, "Санкт-Петербургский государственный университет", size=Pt(14), color=DK1, space_before=Pt(4))
    add_paragraph(tf2, "аэрокосмического приборостроения (ГУАП)", size=Pt(14), color=DK1, space_before=Pt(0))

    txBox3 = add_textbox(slide, Emu(8152598), Emu(5993952), Emu(3488540), Emu(387798))
    set_text(txBox3.text_frame, "2025", size=Pt(18), color=DK1, align=PP_ALIGN.RIGHT)

    # ─── SLIDE 2: Введение — что такое оркестрация ───────────────────
    slide = prs.slides.add_slide(layout)
    clear_slide(slide)
    add_title(slide, "Оркестрация в программной инженерии")

    # Definition block
    add_rounded_rect(slide, LEFT_MARGIN, Emu(1000000), Emu(10800000), Emu(800000),
                     ACCENT1,
                     "Оркестрация — автоматическая координация взаимосвязанных процессов, сервисов и данных для достижения целостного результата",
                     WHITE, Pt(16), True)

    # Left: analogy
    txBox = add_textbox(slide, LEFT_MARGIN, Emu(2100000), Emu(5200000), Emu(4200000))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_text(tf, "Аналогия с дирижёром", size=Pt(17), bold=True, color=DARK_BLUE)
    add_paragraph(tf, "Как дирижёр оркестра не играет на инструментах сам, а координирует музыкантов — так и оркестратор не выполняет задачи, а управляет их порядком, зависимостями и потоками данных.", size=Pt(14), color=DK1, space_before=Pt(10))

    add_paragraph(tf, "Ключевые свойства:", size=Pt(16), bold=True, color=ACCENT1, space_before=Pt(18))
    add_bullet(tf, "Централизованное управление потоком", size=Pt(14))
    add_bullet(tf, "Определённый порядок выполнения шагов", size=Pt(14))
    add_bullet(tf, "Обработка ошибок и откатов", size=Pt(14))
    add_bullet(tf, "Мониторинг состояния всего конвейера", size=Pt(14))

    # Right: diagram — orchestrator in center
    cx = Emu(8900000)
    cy = Emu(3400000)
    add_circle(slide, cx - Emu(500000), cy - Emu(500000), Emu(1000000), ACCENT1, "Оркестратор", WHITE, Pt(11))

    satellites = [
        ("Сервис A", Emu(-1800000), Emu(-1200000), ACCENT3),
        ("Сервис B", Emu(1800000), Emu(-1200000), ACCENT4),
        ("Данные", Emu(-1800000), Emu(1200000), ACCENT5),
        ("Результат", Emu(1800000), Emu(1200000), ACCENT6),
    ]
    for label, dx, dy, clr in satellites:
        sx = cx + dx - Emu(400000)
        sy = cy + dy - Emu(300000)
        add_rounded_rect(slide, sx, sy, Emu(1300000), Emu(600000), clr, label, WHITE, Pt(11), True)
        # Arrow from center to satellite
        add_arrow(slide, cx, cy, sx + Emu(650000), sy + Emu(300000), clr, Pt(2))

    # ─── SLIDE 3: Оркестрация vs хореография vs автоматизация ────────
    slide = prs.slides.add_slide(layout)
    clear_slide(slide)
    add_title(slide, "Оркестрация, хореография и автоматизация")

    col_w = Emu(3500000)
    col_h = Emu(4800000)
    gap = Emu(200000)
    start_x = Emu(400000)
    start_y = Emu(1100000)

    comparisons = [
        ("Автоматизация", ACCENT5,
         ["Отдельный процесс\nвыполняется без участия\nчеловека",
          "Нет координации\nмежду процессами",
          "Пример: скрипт сборки,\ncron-задача"]),
        ("Оркестрация", ACCENT1,
         ["Центральный контроллер\nуправляет всеми\nучастниками",
          "Знает весь workflow\nцеликом, принимает\nрешения",
          "Пример: Kubernetes,\nApache Airflow, CI/CD"]),
        ("Хореография", ACCENT4,
         ["Участники реагируют\nна события\nсамостоятельно",
          "Нет единого центра,\nкаждый знает\nсвою роль",
          "Пример: микросервисы\nс event bus"]),
    ]

    for i, (title, clr, items) in enumerate(comparisons):
        x = start_x + Emu(i * 3900000)

        # Header
        add_rect(slide, x, start_y, col_w, Emu(550000), clr, title, WHITE, Pt(18), True)

        # Items
        for j, item_text in enumerate(items):
            iy = start_y + Emu(700000) + Emu(j * 1350000)
            add_rounded_rect(slide, x + Emu(50000), iy, col_w - Emu(100000), Emu(1150000),
                             LIGHT_BG, item_text, DK1, Pt(13))

    # Bottom note
    txb = add_textbox(slide, LEFT_MARGIN, Emu(6200000), Emu(10800000), Emu(400000))
    set_text(txb.text_frame, "Оркестрация сочетает контроль автоматизации с масштабом координации хореографии",
             size=Pt(14), italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # ─── SLIDE 4: Области применения ─────────────────────────────────
    slide = prs.slides.add_slide(layout)
    clear_slide(slide)
    add_title(slide, "Области применения оркестрации")

    domains = [
        ("Контейнеры\nи облако", "Kubernetes, Docker\nSwarm, Nomad", ACCENT1),
        ("CI / CD", "Jenkins, GitLab CI,\nGitHub Actions", ACCENT3),
        ("Data\npipelines", "Apache Airflow,\nLuigi, Prefect", ACCENT4),
        ("Бизнес-\nпроцессы", "Camunda, Temporal,\nApache Camel", ACCENT5),
        ("IoT и\nвстраиваемые", "Node-RED,\nAWS IoT Greengrass", ACCENT6),
        ("Документо-\nоборот", "Генерация отчётов,\nмиграции данных", ACCENT2),
    ]

    for i, (name, examples, clr) in enumerate(domains):
        col = i % 3
        row = i // 3
        x = Emu(400000 + col * 3900000)
        y = Emu(1150000 + row * 2800000)

        # Card
        add_rounded_rect(slide, x, y, Emu(3500000), Emu(2400000), clr)

        # Name
        txb = add_textbox(slide, x + Emu(150000), y + Emu(200000), Emu(3200000), Emu(800000))
        set_text(txb.text_frame, name, size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Examples
        txb2 = add_textbox(slide, x + Emu(150000), y + Emu(1200000), Emu(3200000), Emu(1000000))
        txb2.text_frame.word_wrap = True
        set_text(txb2.text_frame, examples, size=Pt(14), color=WHITE, align=PP_ALIGN.CENTER)

    # ─── SLIDE 5: Архитектурные паттерны ─────────────────────────────
    slide = prs.slides.add_slide(layout)
    clear_slide(slide)
    add_title(slide, "Архитектурные паттерны оркестрации")

    # Left column — patterns list
    txBox = add_textbox(slide, LEFT_MARGIN, TOP_CONTENT, Emu(5400000), Emu(5200000))
    tf = txBox.text_frame
    tf.word_wrap = True

    patterns = [
        ("Центральный контроллер (Mediator)", "Один компонент знает весь workflow и вызывает участников последовательно"),
        ("Конвейер (Pipeline / Chain)", "Данные проходят через цепочку обработчиков, каждый выполняет свою трансформацию"),
        ("Saga-паттерн", "Длительная транзакция из шагов с компенсирующими действиями при ошибке"),
        ("DAG (направленный ациклический граф)", "Задачи образуют граф зависимостей; выполняются параллельно где возможно"),
    ]

    for i, (name, desc) in enumerate(patterns):
        if i == 0:
            set_text(tf, name, size=Pt(15), bold=True, color=ACCENT1)
        else:
            add_paragraph(tf, name, size=Pt(15), bold=True, color=ACCENT1, space_before=Pt(16))
        add_paragraph(tf, desc, size=Pt(13), color=DK1, space_before=Pt(4))

    # Right — visual: pipeline diagram
    txb_r = add_textbox(slide, Emu(6300000), Emu(1000000), Emu(5000000), Emu(400000))
    set_text(txb_r.text_frame, "Паттерн «Конвейер»", size=Pt(16), bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)

    pipe_steps = [
        ("Вход", ACCENT3),
        ("Парсинг", ACCENT4),
        ("Валидация", ACCENT1),
        ("Трансформация", ACCENT5),
        ("Выход", ACCENT6),
    ]
    for i, (label, clr) in enumerate(pipe_steps):
        y = Emu(1500000 + i * 900000)
        add_rounded_rect(slide, Emu(6800000), y, Emu(4000000), Emu(580000),
                         clr, label, WHITE, Pt(14), True)
        if i < len(pipe_steps) - 1:
            add_arrow(slide, Emu(8800000), y + Emu(580000), Emu(8800000), y + Emu(900000), clr, Pt(2))

    # Bottom — DAG mini-diagram
    txb_dag = add_textbox(slide, Emu(6300000), Emu(5800000), Emu(5000000), Emu(350000))
    set_text(txb_dag.text_frame, "Паттерн «DAG»: параллельное выполнение независимых задач",
             size=Pt(12), italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # ─── SLIDE 6: Ключевые компоненты оркестратора ───────────────────
    slide = prs.slides.add_slide(layout)
    clear_slide(slide)
    add_title(slide, "Ключевые компоненты программного оркестратора")

    components = [
        ("Планировщик\n(Scheduler)", "Определяет порядок\nи время запуска задач", ACCENT1),
        ("Менеджер\nсостояний", "Хранит текущий этап\nworkflow и контекст", ACCENT3),
        ("Адаптеры\nисточников", "Получают данные из\nвнешних систем и API", ACCENT4),
        ("Движок\nтрансформаций", "Преобразует данные\nмежду этапами", ACCENT5),
        ("Обработчик\nошибок", "Retry, fallback,\nкомпенсация, логирование", ACCENT2),
        ("Интерфейс\nмониторинга", "Визуализация прогресса\nи метрик выполнения", ACCENT6),
    ]

    for i, (name, desc, clr) in enumerate(components):
        col = i % 3
        row = i // 3
        x = Emu(400000 + col * 3900000)
        y = Emu(1100000 + row * 2700000)

        # Icon circle
        add_circle(slide, x, y + Emu(200000), Emu(700000), clr, str(i + 1), WHITE, Pt(24))

        # Text
        txb = add_textbox(slide, x + Emu(850000), y + Emu(100000), Emu(2700000), Emu(600000))
        set_text(txb.text_frame, name.replace('\n', ' '), size=Pt(15), bold=True, color=DK1)
        txb2 = add_textbox(slide, x + Emu(850000), y + Emu(700000), Emu(2700000), Emu(1400000))
        txb2.text_frame.word_wrap = True
        set_text(txb2.text_frame, desc, size=Pt(13), color=GRAY)

    # ─── SLIDE 7: Преимущества и вызовы ──────────────────────────────
    slide = prs.slides.add_slide(layout)
    clear_slide(slide)
    add_title(slide, "Преимущества и вызовы оркестрации")

    # Left: advantages
    add_rect(slide, LEFT_MARGIN, Emu(1000000), Emu(5200000), Emu(550000),
             ACCENT6, "Преимущества", WHITE, Pt(18), True)

    txBox = add_textbox(slide, LEFT_MARGIN, Emu(1700000), Emu(5200000), Emu(4500000))
    tf = txBox.text_frame
    tf.word_wrap = True

    advantages = [
        ("Прозрачность", "Весь процесс виден и предсказуем"),
        ("Повторяемость", "Одинаковый результат при одинаковых входных данных"),
        ("Масштабируемость", "Легко добавлять новые этапы и источники"),
        ("Обработка ошибок", "Централизованный retry, fallback, компенсация"),
        ("Наблюдаемость", "Логирование и мониторинг каждого шага"),
    ]

    for i, (term, desc) in enumerate(advantages):
        if i == 0:
            set_text(tf, f"{term} — {desc}", size=Pt(14), color=DK1)
        else:
            add_bullet(tf, f"{term} — {desc}", size=Pt(14))

    # Right: challenges
    add_rect(slide, Emu(6300000), Emu(1000000), Emu(5200000), Emu(550000),
             ACCENT2, "Вызовы", WHITE, Pt(18), True)

    txBox2 = add_textbox(slide, Emu(6300000), Emu(1700000), Emu(5200000), Emu(4500000))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    challenges = [
        ("Единая точка отказа", "Сбой оркестратора останавливает всё"),
        ("Сложность отладки", "Многошаговые конвейеры трудно трассировать"),
        ("Связанность", "Оркестратор должен знать обо всех участниках"),
        ("Производительность", "Централизация может стать узким местом"),
    ]

    for i, (term, desc) in enumerate(challenges):
        if i == 0:
            set_text(tf2, f"{term} — {desc}", size=Pt(14), color=DK1)
        else:
            add_bullet(tf2, f"{term} — {desc}", size=Pt(14))

    # Bottom: balance note
    txb_note = add_textbox(slide, LEFT_MARGIN, Emu(6100000), Emu(10800000), Emu(500000))
    set_text(txb_note.text_frame,
             "Баланс: оркестрация оправдана при сложных workflow с множеством зависимостей и строгими требованиями к порядку выполнения",
             size=Pt(13), italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # ─── SLIDE 8: Технологический ландшафт ───────────────────────────
    slide = prs.slides.add_slide(layout)
    clear_slide(slide)
    add_title(slide, "Технологический ландшафт")

    # Table-like grid
    categories = [
        ("Контейнерная оркестрация", [("Kubernetes", "Google → CNCF"), ("Docker Swarm", "Docker Inc"), ("Nomad", "HashiCorp")], ACCENT1),
        ("Workflow / Data", [("Apache Airflow", "Airbnb → Apache"), ("Prefect", "Prefect.io"), ("Temporal", "Temporal.io")], ACCENT4),
        ("CI/CD", [("Jenkins", "Open Source"), ("GitLab CI", "GitLab"), ("GitHub Actions", "GitHub")], ACCENT3),
        ("Бизнес-процессы", [("Camunda", "Camunda"), ("Apache Camel", "Apache"), ("n8n", "Open Source")], ACCENT5),
    ]

    for i, (cat_name, tools, clr) in enumerate(categories):
        y = Emu(1100000 + i * 1300000)

        # Category label
        add_rect(slide, LEFT_MARGIN, y, Emu(2800000), Emu(1050000), clr, cat_name, WHITE, Pt(14), True)

        # Tool cards
        for j, (tool, origin) in enumerate(tools):
            tx = Emu(3600000 + j * 2900000)
            add_rounded_rect(slide, tx, y + Emu(50000), Emu(2600000), Emu(950000), LIGHT_BG)
            txb = add_textbox(slide, tx + Emu(100000), y + Emu(130000), Emu(2400000), Emu(400000))
            set_text(txb.text_frame, tool, size=Pt(16), bold=True, color=DK1, align=PP_ALIGN.CENTER)
            txb2 = add_textbox(slide, tx + Emu(100000), y + Emu(550000), Emu(2400000), Emu(350000))
            set_text(txb2.text_frame, origin, size=Pt(11), color=GRAY, align=PP_ALIGN.CENTER)

    # ─── SLIDE 9: Принципы проектирования ────────────────────────────
    slide = prs.slides.add_slide(layout)
    clear_slide(slide)
    add_title(slide, "Принципы проектирования оркестратора")

    principles = [
        ("Идемпотентность", "Повторный запуск шага даёт тот же результат.\nПозволяет безопасно перезапускать при сбоях.", ACCENT1),
        ("Слабая связанность", "Участники не знают друг о друге — только оркестратор\nзнает полный workflow. Упрощает замену компонентов.", ACCENT3),
        ("Наблюдаемость", "Каждый шаг логирует вход, выход и метрики.\nДаёт возможность диагностики и оптимизации.", ACCENT4),
        ("Декларативность", "Описание workflow отделено от реализации шагов.\nDAG/YAML/JSON вместо императивного кода.", ACCENT5),
        ("Отказоустойчивость", "Retry с backoff, circuit breaker, компенсирующие\nтранзакции. Система продолжает работу при частичных сбоях.", ACCENT2),
    ]

    for i, (name, desc, clr) in enumerate(principles):
        y = Emu(1050000 + i * 1100000)
        # Color bar
        add_rect(slide, LEFT_MARGIN, y, Emu(80000), Emu(900000), clr)
        # Title
        txb = add_textbox(slide, Emu(800000), y + Emu(50000), Emu(2800000), Emu(400000))
        set_text(txb.text_frame, name, size=Pt(17), bold=True, color=clr)
        # Description
        txb2 = add_textbox(slide, Emu(3800000), y + Emu(50000), Emu(7500000), Emu(850000))
        txb2.text_frame.word_wrap = True
        set_text(txb2.text_frame, desc, size=Pt(13), color=DK1)

    # ─── SLIDE 10: Пример — ПАС-помощник ────────────────────────────
    slide = prs.slides.add_slide(layout)
    clear_slide(slide)
    add_title(slide, "Практический пример: оркестратор паспортов ПАС")

    # Subtitle
    txb_sub = add_textbox(slide, LEFT_MARGIN, Emu(900000), Emu(10800000), Emu(350000))
    set_text(txb_sub.text_frame,
             "Применение принципов оркестрации при формировании паспортов аппаратно-программных средств",
             size=Pt(14), italic=True, color=GRAY)

    # Pipeline flow
    steps = [
        ("Ввод\nпараметров", ACCENT1),
        ("Парсинг\npci -vvv", ACCENT3),
        ("Идентификация\nустройств", ACCENT4),
        ("Привязка\nк подсистемам", ACCENT5),
        ("Генерация\nотчёта .txt", ACCENT6),
    ]

    for i, (label, clr) in enumerate(steps):
        x = Emu(200000 + i * 2350000)
        add_chevron(slide, x, Emu(1400000), Emu(2150000), Emu(800000), clr, label, WHITE, Pt(12))

    # Key orchestration aspects applied
    txBox = add_textbox(slide, LEFT_MARGIN, Emu(2500000), Emu(5400000), Emu(4000000))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_text(tf, "Реализованные принципы:", size=Pt(16), bold=True, color=ACCENT1)
    add_bullet(tf, "Конвейер: данные проходят через 5 последовательных этапов обработки", size=Pt(13))
    add_bullet(tf, "Центральный контроллер: AppController координирует все модули", size=Pt(13))
    add_bullet(tf, "Адаптеры: PCI-парсер, база pci.ids, пресеты подсистем", size=Pt(13))
    add_bullet(tf, "Менеджер состояний: автосохранение прогресса в .progress.json", size=Pt(13))
    add_bullet(tf, "Валидация: контроль формата на каждом этапе", size=Pt(13))

    # Right: tech stack mini
    txBox2 = add_textbox(slide, Emu(6300000), Emu(2500000), Emu(5000000), Emu(400000))
    set_text(txBox2.text_frame, "Стек реализации:", size=Pt(16), bold=True, color=ACCENT4)

    stack_items = [
        ("C++17 + Qt 6", ACCENT1),
        ("QML (Material.Dark)", ACCENT3),
        ("pci.ids (>20 000 записей)", ACCENT4),
        ("JSON-пресеты подсистем", ACCENT5),
        ("Jenkins PDF-генератор", ACCENT2),
    ]
    for i, (item, clr) in enumerate(stack_items):
        y = Emu(3050000 + i * 650000)
        add_rounded_rect(slide, Emu(6300000), y, Emu(5000000), Emu(500000), clr, item, WHITE, Pt(13), True)

    # ─── SLIDE 11: Метрики и результаты ──────────────────────────────
    slide = prs.slides.add_slide(layout)
    clear_slide(slide)
    add_title(slide, "Эффективность оркестрации: до и после")

    # Before / After comparison
    add_rect(slide, LEFT_MARGIN, Emu(1050000), Emu(5200000), Emu(500000),
             ACCENT2, "Без оркестрации (вручную)", WHITE, Pt(16), True)

    manual_items = [
        ("Сбор данных PCI", "~30 мин", 0.8),
        ("Заполнение таблиц", "~60 мин", 1.0),
        ("Форматирование разметки", "~45 мин", 0.9),
        ("Проверка на ошибки", "~20 мин", 0.5),
    ]
    for i, (label, time, pct) in enumerate(manual_items):
        y = Emu(1700000 + i * 500000)
        txb = add_textbox(slide, Emu(600000), y, Emu(2500000), Emu(400000))
        set_text(txb.text_frame, label, size=Pt(12), color=DK1)
        bar_w = int(2200000 * pct)
        add_rect(slide, Emu(3200000), y + Emu(50000), Emu(bar_w), Emu(300000), ACCENT2)
        txb_t = add_textbox(slide, Emu(3200000) + Emu(bar_w) + Emu(80000), y + Emu(30000), Emu(1000000), Emu(350000))
        set_text(txb_t.text_frame, time, size=Pt(11), bold=True, color=ACCENT2)

    add_rect(slide, Emu(6300000), Emu(1050000), Emu(5200000), Emu(500000),
             ACCENT6, "С программным оркестратором", WHITE, Pt(16), True)

    auto_items = [
        ("Вставка pci -vvv", "~1 мин", 0.10),
        ("Автозаполнение таблиц", "~3 мин", 0.15),
        ("Генерация разметки", "мгновенно", 0.05),
        ("Автовалидация", "мгновенно", 0.05),
    ]
    for i, (label, time, pct) in enumerate(auto_items):
        y = Emu(1700000 + i * 500000)
        txb = add_textbox(slide, Emu(6400000), y, Emu(2500000), Emu(400000))
        set_text(txb.text_frame, label, size=Pt(12), color=DK1)
        bar_w = int(2200000 * pct)
        add_rect(slide, Emu(9000000), y + Emu(50000), Emu(bar_w), Emu(300000), ACCENT6)
        txb_t = add_textbox(slide, Emu(9000000) + Emu(bar_w) + Emu(80000), y + Emu(30000), Emu(1000000), Emu(350000))
        set_text(txb_t.text_frame, time, size=Pt(11), bold=True, color=ACCENT6)

    # Summary line
    add_rect(slide, LEFT_MARGIN, Emu(4000000), Emu(10800000), Emu(60000), ACCENT1)

    # Bottom: general conclusions
    txBox = add_textbox(slide, LEFT_MARGIN, Emu(4200000), Emu(10800000), Emu(2600000))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_text(tf, "Общие показатели эффективности оркестрации:", size=Pt(16), bold=True, color=ACCENT1)
    add_bullet(tf, "Сокращение ручного труда в 5–10 раз за счёт автоматизации рутинных этапов", size=Pt(14))
    add_bullet(tf, "Устранение класса ошибок форматирования — валидация на каждом шаге конвейера", size=Pt(14))
    add_bullet(tf, "Повторяемость: один и тот же процесс гарантирует одинаковый результат", size=Pt(14))
    add_bullet(tf, "Прозрачность: оператор видит текущий этап и может вмешаться при необходимости", size=Pt(14))

    # ─── SLIDE 12: Выводы ───────────────────────────────────────────
    slide = prs.slides.add_slide(layout)
    clear_slide(slide)
    add_title(slide, "Выводы")

    conclusions = [
        ("Оркестрация — зрелый подход", "к координации сложных многоэтапных процессов, применимый от облачных систем до генерации документов", ACCENT1),
        ("Ключевые паттерны", "— конвейер, медиатор, saga, DAG — покрывают большинство сценариев координации", ACCENT3),
        ("Принципы проектирования", "— идемпотентность, слабая связанность, наблюдаемость — обеспечивают надёжность", ACCENT4),
        ("Практическая реализация", "подтверждает: даже небольшой оркестратор значительно сокращает трудозатраты и ошибки", ACCENT5),
        ("Перспективы развития", "— интеграция с внешними API, расширение адаптеров, декларативные описания workflow", ACCENT6),
    ]

    for i, (term, desc, clr) in enumerate(conclusions):
        y = Emu(1100000 + i * 1050000)
        add_circle(slide, LEFT_MARGIN, y + Emu(50000), Emu(600000), clr, str(i + 1), WHITE, Pt(22))
        txb = add_textbox(slide, Emu(1350000), y + Emu(50000), Emu(10000000), Emu(900000))
        txb.text_frame.word_wrap = True
        p = set_text(txb.text_frame, term, size=Pt(15), bold=True, color=clr)
        run2 = p.add_run()
        run2.text = " " + desc
        run2.font.size = Pt(15)
        run2.font.color.rgb = DK1
        run2.font.name = FONT

    # ─── SLIDE 13: Спасибо ──────────────────────────────────────────
    slide = prs.slides.add_slide(layout)
    clear_slide(slide)

    txBox = add_textbox(slide, Emu(550864), Emu(1800000), Emu(10800000), Emu(800000))
    set_text(txBox.text_frame, "Спасибо за внимание!", size=Pt(44), color=ACCENT1, align=PP_ALIGN.CENTER)

    txBox2 = add_textbox(slide, Emu(550864), Emu(3200000), Emu(10800000), Emu(1800000))
    tf2 = txBox2.text_frame
    set_text(tf2, "Соколов Илья Дмитриевич", size=Pt(22), color=DK1, align=PP_ALIGN.CENTER)
    add_paragraph(tf2, "ГУАП, Санкт-Петербург", size=Pt(18), color=DK1, align=PP_ALIGN.CENTER, space_before=Pt(8))
    add_paragraph(tf2, "Готов ответить на ваши вопросы", size=Pt(16), color=ACCENT1,
                  align=PP_ALIGN.CENTER, space_before=Pt(24), italic=True)

    txBox3 = add_textbox(slide, Emu(8152598), Emu(5993952), Emu(3488540), Emu(387798))
    set_text(txBox3.text_frame, "2025", size=Pt(18), color=DK1, align=PP_ALIGN.RIGHT)

    # ── Save ─────────────────────────────────────────────────────────
    output_path = os.path.join(os.path.dirname(__file__), 'presentation.pptx')
    prs.save(output_path)
    print(f"Saved: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
